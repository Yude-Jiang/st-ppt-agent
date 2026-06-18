"""ST brand palette, helpers, and verified builder functions.

Palette constants and builder functions (add_message_bar, add_shaded_box,
add_section_slide, add_cards_row, add_comparison, add_process_flow) are
copied verbatim from st-ppt-brand/references/pptx-implementation.md.
Do NOT modify their logic — they are brand-validated.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Palette (official ST names: Green Vogue / Gold / Picton Blue / White)
# ---------------------------------------------------------------------------
ST_DARK_BLUE = RGBColor(0x03, 0x23, 0x4B)
ST_YELLOW    = RGBColor(0xFF, 0xD2, 0x00)
ST_LIGHT_BLUE = RGBColor(0x3C, 0xB4, 0xE6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_1       = RGBColor(0xEE, 0xEF, 0xF1)
GRAY_2       = RGBColor(0xDB, 0xDE, 0xE1)
GRAY_3       = RGBColor(0xC0, 0xC8, 0xD2)

RAMP = [
    RGBColor(0x03, 0x23, 0x4B),  # 1 ST Dark Blue
    RGBColor(0x42, 0x59, 0x78),  # 2 slate
    RGBColor(0x80, 0x91, 0xA5),  # 3 medium blue-gray
    RGBColor(0xC0, 0xC9, 0xCE),  # 4 light blue-gray
]

FONT = "Arial"

RECT = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE
OVAL = 9   # MSO_AUTO_SHAPE_TYPE.OVAL


def ramp_text(step):
    return WHITE if step < 2 else ST_DARK_BLUE


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def no_autofit(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def style_runs(tf, color, size_pt, bold=False, font=FONT):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.color.rgb = color


# ---------------------------------------------------------------------------
# Builder functions — verbatim from pptx-implementation.md
# ---------------------------------------------------------------------------

def add_message_bar(slide, text, fill_color=ST_LIGHT_BLUE):
    bar = slide.shapes.add_shape(
        RECT,
        left=Inches(0), top=Inches(1.43),
        width=Inches(9.8), height=Inches(0.84))
    fill(bar, fill_color)
    tf = bar.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    txt = WHITE if fill_color == ST_DARK_BLUE else ST_DARK_BLUE
    style_runs(tf, txt, 20, bold=True)
    return bar


def add_shaded_box(slide, x, y, w, h, lines, shade=GRAY_1, heading=None):
    box = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(box, shade)
    tf = box.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    if heading:
        tf.text = heading
        style_runs(tf, ST_DARK_BLUE, 14, bold=True)
        first = False
    else:
        first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + ln
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(14)
            r.font.color.rgb = ST_DARK_BLUE
    return box


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = ST_DARK_BLUE
    bar = slide.shapes.add_shape(RECT, Inches(2.2), Inches(0.9), Inches(9.0), Inches(1.2))
    fill(bar, ST_YELLOW)
    tf = bar.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    tf.text = title
    style_runs(tf, ST_DARK_BLUE, 32, bold=True)
    return slide


def add_cards_row(slide, cards, top=2.2, height=4.6, gap=0.3,
                  left_margin=0.4, header="yellow", with_image=True):
    """cards = [{"title": str, "bullets": [...], "img": path|None}]"""
    n = len(cards)
    total_w = 13.333 - 2 * left_margin
    w = (total_w - gap * (n - 1)) / n
    hdr_h = 0.7
    img_h = 2.2 if with_image else 0
    for i, c in enumerate(cards):
        x = left_margin + i * (w + gap)
        if header == "yellow":
            hfill, htext = ST_YELLOW, ST_DARK_BLUE
        else:
            hfill, htext = RAMP[min(i, 3)], ramp_text(min(i, 3))
        hb = slide.shapes.add_shape(RECT, Inches(x), Inches(top),
                                    Inches(w), Inches(hdr_h))
        fill(hb, hfill)
        tf = hb.text_frame
        no_autofit(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.text = c["title"]
        style_runs(tf, htext, 14, bold=True)
        y = top + hdr_h
        if with_image and c.get("img"):
            slide.shapes.add_picture(c["img"], Inches(x), Inches(y),
                                     width=Inches(w), height=Inches(img_h))
            y += img_h
        gb = slide.shapes.add_shape(RECT, Inches(x), Inches(y),
                                    Inches(w), Inches(top + height - y))
        fill(gb, GRAY_1)
        tf = gb.text_frame
        no_autofit(tf)
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.12)
        first = True
        for b in c.get("bullets", []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "• " + b
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(14)
                r.font.color.rgb = ST_DARK_BLUE


def add_comparison(slide, left, right, top=1.7):
    """left/right = {"name", "desc", "img", "bullets"}"""
    for col, data, hdr in [(0, left, ST_DARK_BLUE), (1, right, ST_LIGHT_BLUE)]:
        x = 0.5 + col * 6.4
        w = 5.9
        hb = slide.shapes.add_shape(RECT, Inches(x), Inches(top),
                                    Inches(w), Inches(0.9))
        fill(hb, hdr)
        tf = hb.text_frame
        no_autofit(tf)
        tf.margin_left = Inches(0.2)
        tf.text = data["name"]
        style_runs(tf, WHITE, 16, bold=True)
        p = tf.add_paragraph()
        p.text = data.get("desc", "")
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(13)
            r.font.color.rgb = WHITE
        if data.get("img"):
            slide.shapes.add_picture(data["img"], Inches(x), Inches(top + 0.9),
                                     width=Inches(w), height=Inches(2.6))
        gb = slide.shapes.add_shape(RECT, Inches(x), Inches(top + 3.5),
                                    Inches(w), Inches(1.8))
        fill(gb, GRAY_1)
        tf = gb.text_frame
        no_autofit(tf)
        tf.margin_left = Inches(0.2)
        first = True
        for b in data.get("bullets", []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "• " + b
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(14)
                r.font.color.rgb = ST_DARK_BLUE


def add_process_flow(slide, steps, cy=4.1, d=1.9):
    """steps = [{"label": circle text, "step": "Step N", "caption": text}]"""
    n = len(steps)
    span = 13.333 - 1.2
    gap = (span - d * n) / (n - 1) if n > 1 else 0
    for i, s in enumerate(steps):
        x = 0.6 + i * (d + gap)
        circ = slide.shapes.add_shape(OVAL, Inches(x), Inches(cy - d / 2),
                                      Inches(d), Inches(d))
        col = RAMP[0] if i % 2 == 0 else RAMP[2]
        fill(circ, col)
        tf = circ.text_frame
        no_autofit(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.text = s["label"]
        style_runs(tf, ramp_text(0 if i % 2 == 0 else 2), 12, bold=True)
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        my = cy - d / 2 - 0.55 if i % 2 == 0 else cy + d / 2 + 0.35
        dot = slide.shapes.add_shape(OVAL, Inches(x + d / 2 - 0.12), Inches(my),
                                     Inches(0.24), Inches(0.24))
        fill(dot, ST_YELLOW)
