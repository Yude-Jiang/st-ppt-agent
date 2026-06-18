"""Load and delegate to the real st-ppt-brand skill builder."""
from __future__ import annotations

import importlib.util
import io
import logging
import os
import sys
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation

logger = logging.getLogger(__name__)

VENDOR_DIR = Path(__file__).resolve().parent.parent / "vendor" / "st_ppt_brand"

_BUILD_FN: Callable[[list[dict[str, Any]]], bytes] | None = None
_FIELD_SPECS: dict[str, str] | None = None

# st-ppt-brand skill may expose per-archetype add_* helpers (see exploration doc).
_ARCHETYPE_ADDER_NAMES: dict[str, list[str]] = {
    "title-slide": ["add_title_slide"],
    "title-bullets": ["add_title_bullets", "add_bullets"],
    "title-image": ["add_title_image"],
    "two-column": ["add_two_column"],
    "three-column": ["add_three_column"],
    "product-comparison-2up": ["add_comparison", "add_product_comparison_2up"],
    "process-flow": ["add_process_flow"],
    "cards-row": ["add_cards_row"],
    "quote-highlight": ["add_quote_highlight"],
    "content-placeholder": ["add_content_placeholder"],
    "section-divider": ["add_section_divider"],
}

_BUILDER_MODULE_NAMES = ("build_deck", "builder", "render", "ppt_builder")
_BUILD_FN_NAMES = ("build_pptx", "build_deck", "render_deck", "render_pptx")


class SkillNotFoundError(RuntimeError):
    """Raised when st-ppt-brand skill cannot be located."""


def skill_root() -> Path | None:
    """Return configured skill directory, or None if not available."""
    env_path = os.getenv("ST_PPT_BRAND_SKILL_PATH", "").strip()
    if env_path:
        path = Path(env_path).expanduser().resolve()
        if path.is_dir():
            return path

    if VENDOR_DIR.is_dir() and _has_builder_files(VENDOR_DIR):
        return VENDOR_DIR

    return None


def _has_builder_files(directory: Path) -> bool:
    markers = (
        "build_deck.py",
        "builder.py",
        "render.py",
        "ppt_builder.py",
        "__init__.py",
    )
    if any((directory / name).is_file() for name in markers):
        return True
    scripts = directory / "scripts"
    return scripts.is_dir() and any(scripts.glob("*.py"))


def _load_module_from_path(module_name: str, file_path: Path):
    spec = importlib.util.spec_from_file_location(module_name, file_path)
    if spec is None or spec.loader is None:
        raise ImportError(f"Cannot load module from {file_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    return module


def _import_skill_module(root: Path):
    """Import the primary builder module from a skill directory."""
    for name in _BUILDER_MODULE_NAMES:
        candidate = root / f"{name}.py"
        if candidate.is_file():
            return _load_module_from_path(f"st_ppt_brand_{name}", candidate)

    init_file = root / "__init__.py"
    if init_file.is_file():
        return _load_module_from_path("st_ppt_brand_pkg", init_file)

    scripts = root / "scripts"
    if scripts.is_dir():
        for name in _BUILDER_MODULE_NAMES:
            candidate = scripts / f"{name}.py"
            if candidate.is_file():
                return _load_module_from_path(f"st_ppt_brand_scripts_{name}", candidate)
        py_files = sorted(scripts.glob("*.py"))
        if len(py_files) == 1:
            only = py_files[0]
            return _load_module_from_path(f"st_ppt_brand_scripts_{only.stem}", only)

    raise SkillNotFoundError(
        f"No builder module found under {root}. Expected build_deck.py, builder.py, "
        "or scripts/*.py exporting build_pptx()."
    )


def _resolve_callable(module, names: tuple[str, ...]) -> Callable[..., Any] | None:
    for name in names:
        fn = getattr(module, name, None)
        if callable(fn):
            return fn
    return None


def _build_from_adders(module, items: list[dict[str, Any]]) -> bytes:
    """Compose deck via add_* helpers when skill does not expose build_pptx()."""
    new_prs = _resolve_callable(module, ("new_presentation", "create_presentation", "new_deck"))
    if new_prs is not None:
        prs = new_prs()
    else:
        prs = Presentation()

    add_slide = _resolve_callable(module, ("add_blank_slide", "new_blank_slide"))
    for item in sorted(items, key=lambda x: x.get("order", 0)):
        archetype = item.get("archetype", "title-bullets")
        content_fields = item.get("content_fields", {})
        adder = None
        for candidate in _ARCHETYPE_ADDER_NAMES.get(archetype, []):
            fn = getattr(module, candidate, None)
            if callable(fn):
                adder = fn
                break
        if adder is None:
            raise SkillNotFoundError(
                f"Skill module missing builder for archetype '{archetype}'. "
                f"Tried: {_ARCHETYPE_ADDER_NAMES.get(archetype, [])}"
            )

        if add_slide is not None:
            slide = add_slide(prs)
            adder(slide, content_fields)
        else:
            # Some skills expect (prs, fields) and manage slides internally.
            adder(prs, content_fields)

    save_fn = _resolve_callable(module, ("save_presentation", "presentation_to_bytes", "save_deck"))
    if save_fn is not None:
        result = save_fn(prs)
        if isinstance(result, bytes):
            return result

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _wrap_module_builder(module) -> Callable[[list[dict[str, Any]]], bytes]:
    build_fn = _resolve_callable(module, _BUILD_FN_NAMES)
    if build_fn is not None:
        return build_fn

    renderers = getattr(module, "RENDERERS", None)
    if isinstance(renderers, dict) and renderers:
        def _via_renderers(items: list[dict[str, Any]]) -> bytes:
            prs = Presentation()
            for item in sorted(items, key=lambda x: x.get("order", 0)):
                archetype = item.get("archetype", "title-bullets")
                renderer = renderers.get(archetype)
                if not callable(renderer):
                    raise SkillNotFoundError(f"Skill RENDERERS missing '{archetype}'")
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                renderer(slide, item.get("content_fields", {}))
            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()

        return _via_renderers

    return _build_from_adders


def get_build_pptx() -> Callable[[list[dict[str, Any]]], bytes]:
    """Return cached build_pptx callable from st-ppt-brand skill or legacy fallback."""
    global _BUILD_FN
    if _BUILD_FN is not None:
        return _BUILD_FN

    root = skill_root()
    if root is not None:
        module = _import_skill_module(root)
        _BUILD_FN = _wrap_module_builder(module)
        logger.info("st-ppt-brand skill loaded from %s", root)
        return _BUILD_FN

    if os.getenv("ST_PPT_BRAND_ALLOW_LEGACY", "").lower() in ("1", "true", "yes"):
        from .legacy_renderer import build_pptx as legacy_build

        logger.warning(
            "st-ppt-brand skill not found; using legacy_renderer fallback. "
            "Run scripts/sync-st-ppt-brand.sh or set ST_PPT_BRAND_SKILL_PATH."
        )
        _BUILD_FN = legacy_build
        return _BUILD_FN

    raise SkillNotFoundError(
        "st-ppt-brand skill not found. Set ST_PPT_BRAND_SKILL_PATH, run "
        "scripts/sync-st-ppt-brand.sh, or set ST_PPT_BRAND_ALLOW_LEGACY=1 for dev."
    )


def get_archetype_field_specs() -> dict[str, str]:
    """Return ARCHETYPE_FIELD_SPECS from skill if available, else local fallback."""
    global _FIELD_SPECS
    if _FIELD_SPECS is not None:
        return _FIELD_SPECS

    root = skill_root()
    if root is not None:
        try:
            module = _import_skill_module(root)
            specs = getattr(module, "ARCHETYPE_FIELD_SPECS", None)
            if isinstance(specs, dict) and specs:
                _FIELD_SPECS = {str(k): str(v) for k, v in specs.items()}
                return _FIELD_SPECS
        except Exception as exc:
            logger.warning("Could not load ARCHETYPE_FIELD_SPECS from skill: %s", exc)

    from .archetypes import ARCHETYPE_FIELD_SPECS as fallback

    _FIELD_SPECS = dict(fallback)
    return _FIELD_SPECS


def reset_cache() -> None:
    """Clear cached skill imports (for tests)."""
    global _BUILD_FN, _FIELD_SPECS
    _BUILD_FN = None
    _FIELD_SPECS = None
