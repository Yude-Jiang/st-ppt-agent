"""Layered ecosystem / relationship diagram archetype for ST brand slides."""
from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .brand import (
    ST_DARK_BLUE, ST_YELLOW, ST_LIGHT_BLUE, WHITE, GRAY_1, GRAY_2, GRAY_3,
    FONT, RECT, fill, no_autofit, style_runs,
)

_COLOR_MAP = {
    "dark_blue":  (ST_DARK_BLUE, WHITE),
    "yellow":     (ST_YELLOW, ST_DARK_BLUE),
    "light_blue": (ST_LIGHT_BLUE, WHITE),
    "gray":       (GRAY_1, ST_DARK_BLUE),
}

_LM = 0.4    # left margin, inches
_RW = 12.5   # content row width, inches


def _box(slide, x, y, w, h, bg, fg, title, subtitle=None, size=13,
         bold=True, is_future=False):
    s = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(s, bg)
    if is_future:
        s.line.color.rgb = GRAY_3
        s.line.width = Pt(1.5)
    tf = s.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.text = title
    style_runs(tf, fg, size, bold=bold)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        for r in p2.runs:
            r.font.name = FONT
            r.font.size = Pt(max(size - 2, 9))
            r.font.bold = False
            r.font.color.rgb = fg
        p2.alignment = PP_ALIGN.CENTER


def _arrow(slide, cx, y, color=ST_LIGHT_BLUE):
    s = slide.shapes.add_shape(RECT, Inches(cx - 0.05), Inches(y),
                               Inches(0.1), Inches(0.25))
    fill(s, color)


def render_layered_diagram(slide, cf: dict) -> None:
    """
    Render a layered relationship/ecosystem diagram.

    content_fields:
      title: str
      layers: list of:
        nodes: list[str]          – box labels; use \\n to add a subtitle line
        color: dark_blue | yellow | light_blue | gray
        note: str (optional)      – small caption below the row
        side_nodes: list[str]     – smaller boxes pinned to the right (optional)
        is_future: bool           – dashed border, gray arrows above (optional)
    """
    # Title + yellow accent (same geometry as _slide_title in renderer.py)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.85))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = cf.get("title", "")
    style_runs(tf, ST_DARK_BLUE, 24, bold=True)
    acc = slide.shapes.add_shape(RECT, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.07))
    fill(acc, ST_YELLOW)

    y = 1.3
    ARROW_H = 0.28
    NOTE_H  = 0.22
    GAP      = 0.06
    NODE_GAP = 0.2

    layers = cf.get("layers", [])

    for li, layer in enumerate(layers):
        nodes      = layer.get("nodes", [])
        bg, fg     = _COLOR_MAP.get(layer.get("color", "gray"), (GRAY_1, ST_DARK_BLUE))
        note       = layer.get("note")
        side_nodes = layer.get("side_nodes", [])
        is_future  = layer.get("is_future", False)

        has_sub = any("\n" in n for n in nodes)
        h = 0.95 if has_sub else 0.75

        side_w  = 3.0 if side_nodes else 0.0
        side_gap = 0.2 if side_nodes else 0.0
        main_w  = _RW - side_w - side_gap

        n  = len(nodes)
        nw = (main_w - NODE_GAP * max(n - 1, 0)) / max(n, 1)

        for i, raw in enumerate(nodes):
            parts = raw.split("\n", 1)
            _box(slide,
                 x=_LM + i * (nw + NODE_GAP), y=y, w=nw, h=h,
                 bg=bg, fg=fg,
                 title=parts[0],
                 subtitle=parts[1] if len(parts) > 1 else None,
                 size=14 if n == 1 else 13,
                 is_future=is_future)

        if side_nodes:
            ns = len(side_nodes)
            sh = (h - GAP * max(ns - 1, 0)) / max(ns, 1)
            for si, sn in enumerate(side_nodes):
                _box(slide,
                     x=_LM + main_w + side_gap, y=y + si * (sh + GAP),
                     w=side_w, h=sh,
                     bg=GRAY_2, fg=ST_DARK_BLUE,
                     title=sn, size=10, bold=False)

        if note:
            ntb = slide.shapes.add_textbox(
                Inches(_LM), Inches(y + h + 0.03), Inches(_RW), Inches(NOTE_H))
            ntf = ntb.text_frame
            no_autofit(ntf)
            ntf.text = note
            style_runs(ntf, GRAY_2, 10)
            for p in ntf.paragraphs:
                p.alignment = PP_ALIGN.CENTER

        y += h + (NOTE_H + 0.05 if note else GAP)

        if li < len(layers) - 1:
            next_future = layers[li + 1].get("is_future", False)
            arrow_color = GRAY_2 if next_future else ST_LIGHT_BLUE

            if n <= 1:
                _arrow(slide, _LM + main_w / 2, y, arrow_color)
            else:
                for i in range(n):
                    cx = _LM + i * (nw + NODE_GAP) + nw / 2
                    _arrow(slide, cx, y, arrow_color)

            if next_future:
                fb = slide.shapes.add_textbox(
                    Inches(_LM + main_w / 2 + 0.15), Inches(y + 0.02),
                    Inches(2.0), Inches(0.22))
                ftf = fb.text_frame
                no_autofit(ftf)
                ftf.text = "规划中 →"
                style_runs(ftf, GRAY_2, 10)

            y += ARROW_H + GAP
