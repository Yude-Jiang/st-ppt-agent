"""python-pptx renderer for 11 ST brand archetypes.

Uses ST brand palette and builders from brand.py (sourced from
st-ppt-brand/references/pptx-implementation.md).
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .brand import (
    ST_DARK_BLUE, ST_YELLOW, ST_LIGHT_BLUE, WHITE, GRAY_1, GRAY_2, FONT,
    RECT, fill, no_autofit, style_runs,
    add_message_bar, add_shaded_box, add_cards_row, add_comparison, add_process_flow,
)

W = Inches(13.333)
H = Inches(7.5)


def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_title(slide, text: str) -> None:
    """Title text at top + ST_YELLOW accent bar — used by all content slides."""
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.85))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    style_runs(tf, ST_DARK_BLUE, 24, bold=True)
    acc = slide.shapes.add_shape(RECT, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.07))
    fill(acc, ST_YELLOW)


# ---------------------------------------------------------------------------
# Archetype renderers
# ---------------------------------------------------------------------------

def _render_title_slide(slide, cf: dict) -> None:
    # Dark blue top bar + yellow accent
    top = slide.shapes.add_shape(RECT, 0, 0, W, Inches(0.18))
    fill(top, ST_DARK_BLUE)
    acc = slide.shapes.add_shape(RECT, 0, Inches(0.18), W, Inches(0.18))
    fill(acc, ST_YELLOW)
    # Centered title
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = cf.get("title", "")
    style_runs(tf, ST_DARK_BLUE, 36, bold=True)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    if cf.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.333), Inches(1.0))
        tf = sub.text_frame
        no_autofit(tf)
        tf.text = cf["subtitle"]
        style_runs(tf, GRAY_2, 20)
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
    # Light blue bottom bar
    bot = slide.shapes.add_shape(RECT, 0, Inches(7.2), W, Inches(0.3))
    fill(bot, ST_LIGHT_BLUE)


def _render_title_bullets(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    bullets = cf.get("bullets", [])
    add_shaded_box(slide, 0.4, 1.3, 12.5, 5.9, bullets)


def _render_title_image(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    hint = cf.get("image_hint", "此处放置图片")
    ph = slide.shapes.add_shape(RECT, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.9))
    fill(ph, GRAY_1)
    ph.line.color.rgb = GRAY_2
    tf = ph.text_frame
    no_autofit(tf)
    tf.text = f"[ {hint} ]"
    style_runs(tf, GRAY_2, 16)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER


def _render_two_column(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    col_w = 6.0
    left_bullets = [cf.get("left_content", "")]
    right_bullets = [cf.get("right_content", "")]
    add_shaded_box(slide, 0.4, 1.3, col_w, 5.9, left_bullets,
                   heading=cf.get("left_title", ""))
    add_shaded_box(slide, 6.9, 1.3, col_w, 5.9, right_bullets,
                   heading=cf.get("right_title", ""))


def _render_three_column(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    col_w = 4.0
    cols = [
        (0.4,  cf.get("col1_title", ""), [cf.get("col1_content", "")]),
        (4.7,  cf.get("col2_title", ""), [cf.get("col2_content", "")]),
        (9.0,  cf.get("col3_title", ""), [cf.get("col3_content", "")]),
    ]
    for x, heading, lines in cols:
        add_shaded_box(slide, x, 1.3, col_w, 5.9, lines, heading=heading)


def _render_comparison_2up(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    left_data = {
        "name": cf.get("left_title", ""),
        "desc": "",
        "img": None,
        "bullets": cf.get("left_bullets", []),
    }
    right_data = {
        "name": cf.get("right_title", ""),
        "desc": "",
        "img": None,
        "bullets": cf.get("right_bullets", []),
    }
    add_comparison(slide, left_data, right_data, top=1.5)


def _render_process_flow(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    steps = [
        {"label": s.get("title", ""), "step": "", "caption": s.get("description", "")}
        for s in cf.get("steps", [])
    ]
    if steps:
        add_process_flow(slide, steps)


def _render_cards_row(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    adapted = [
        {"title": c.get("title", ""), "bullets": [c.get("content", "")], "img": None}
        for c in cf.get("cards", [])
    ]
    if adapted:
        add_cards_row(slide, adapted, top=1.5, height=5.7, with_image=False)


def _render_quote_highlight(slide, cf: dict) -> None:
    # Light blue left accent bar
    acc = slide.shapes.add_shape(RECT, 0, Inches(1.5), Inches(0.2), Inches(4.0))
    fill(acc, ST_LIGHT_BLUE)
    # Quote text
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.0), Inches(3.0))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = cf.get("quote", "")
    style_runs(tf, ST_DARK_BLUE, 24, bold=True)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    if cf.get("attribution"):
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.6))
        tf = sub.text_frame
        no_autofit(tf)
        tf.text = f"— {cf['attribution']}"
        style_runs(tf, GRAY_2, 14)
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
    # Bottom yellow accent
    bot = slide.shapes.add_shape(RECT, 0, Inches(7.2), W, Inches(0.3))
    fill(bot, ST_YELLOW)


def _render_content_placeholder(slide, cf: dict) -> None:
    _slide_title(slide, cf.get("title", ""))
    hint = cf.get("placeholder_hint", "此处放置图表")
    ph = slide.shapes.add_shape(RECT, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.9))
    fill(ph, GRAY_1)
    ph.line.color.rgb = GRAY_2
    tf = ph.text_frame
    no_autofit(tf)
    tf.text = f"[ {hint} ]"
    style_runs(tf, GRAY_2, 18)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER


def _render_section_divider(prs: Presentation, cf: dict):
    """Creates and returns a new slide (uses add_section_slide pattern from skill)."""
    from .brand import add_section_slide
    slide = add_section_slide(prs, cf.get("title", ""))
    if cf.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(2.2), Inches(2.4), Inches(9.0), Inches(0.8))
        tf = sub.text_frame
        no_autofit(tf)
        tf.text = cf["subtitle"]
        style_runs(tf, WHITE, 18)
    return slide


# ---------------------------------------------------------------------------
# Dispatch table & entry point
# ---------------------------------------------------------------------------

_RENDERERS = {
    "title-slide": _render_title_slide,
    "title-bullets": _render_title_bullets,
    "title-image": _render_title_image,
    "two-column": _render_two_column,
    "three-column": _render_three_column,
    "product-comparison-2up": _render_comparison_2up,
    "process-flow": _render_process_flow,
    "cards-row": _render_cards_row,
    "quote-highlight": _render_quote_highlight,
    "content-placeholder": _render_content_placeholder,
}


def build_pptx(items: list[dict[str, Any]]) -> bytes:
    """Render a list of SlidePlanItem dicts to .pptx bytes."""
    prs = _prs()
    for item in sorted(items, key=lambda x: x.get("order", 0)):
        archetype = item.get("archetype", "title-bullets")
        cf = item.get("content_fields", {})
        if archetype == "section-divider":
            _render_section_divider(prs, cf)
        else:
            slide = _blank_slide(prs)
            renderer = _RENDERERS.get(archetype, _render_title_bullets)
            renderer(slide, cf)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
