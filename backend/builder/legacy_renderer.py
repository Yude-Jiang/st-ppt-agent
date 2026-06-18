"""Legacy fallback renderer — used only when st-ppt-brand skill is not vendored.

Prefer backend.builder.skill_loader (real st-ppt-brand skill). Set
ST_PPT_BRAND_ALLOW_LEGACY=1 to enable this path in dev/test.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ST brand constants
ST_RED = RGBColor(0xE3, 0x18, 0x37)
ST_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ST_DARK = RGBColor(0x1A, 0x1A, 0x1A)
ST_GREY = RGBColor(0x4D, 0x4D, 0x4D)
ST_LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)

W = Inches(13.333)
H = Inches(7.5)
MSG_H = Inches(0.55)
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.4)
CONTENT_H = H - CONTENT_TOP - MARGIN


def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


def _add_box(slide, left, top, width, height, text: str, font_size: int = 16,
             bold: bool = False, color: RGBColor = ST_DARK,
             bg: RGBColor | None = None, align=PP_ALIGN.LEFT,
             wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg


def _message_bar(slide, title: str) -> None:
    """Red bar at top with white title text."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, MSG_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ST_RED
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {title}"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = ST_WHITE


def _divider_line(slide, top) -> None:
    line = slide.shapes.add_shape(1, MARGIN, top, W - 2 * MARGIN, Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ST_RED
    line.line.fill.background()


def _render_title_slide(slide, cf: dict) -> None:
    # Red bar at top
    _message_bar(slide, "")
    # Big title centered
    _add_box(slide, MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(2.0),
             cf.get("title", ""), font_size=36, bold=True, color=ST_DARK, align=PP_ALIGN.CENTER)
    if cf.get("subtitle"):
        _add_box(slide, MARGIN, Inches(4.2), W - 2 * MARGIN, Inches(1.0),
                 cf["subtitle"], font_size=20, color=ST_GREY, align=PP_ALIGN.CENTER)
    # Red accent line below title
    _divider_line(slide, Inches(4.1))


def _render_title_bullets(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    bullets = cf.get("bullets", [])
    content = "\n".join(f"• {b}" for b in bullets)
    _add_box(slide, MARGIN, CONTENT_TOP, W - 2 * MARGIN, CONTENT_H,
             content, font_size=18, color=ST_DARK)


def _render_title_image(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    hint = cf.get("image_hint", "此处放置图片")
    # Grey placeholder box
    ph = slide.shapes.add_shape(1, MARGIN, CONTENT_TOP, W - 2 * MARGIN, CONTENT_H)
    ph.fill.solid()
    ph.fill.fore_color.rgb = ST_LIGHT_GREY
    ph.line.color.rgb = ST_GREY
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {hint} ]"
    run.font.size = Pt(16)
    run.font.color.rgb = ST_GREY
    run.font.italic = True


def _render_two_column(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    col_w = (W - 3 * MARGIN) / 2
    # Left
    _add_box(slide, MARGIN, CONTENT_TOP, col_w, Inches(0.4),
             cf.get("left_title", ""), font_size=16, bold=True, color=ST_RED)
    _add_box(slide, MARGIN, CONTENT_TOP + Inches(0.5), col_w, CONTENT_H - Inches(0.5),
             cf.get("left_content", ""), font_size=15, color=ST_DARK)
    # Divider
    div = slide.shapes.add_shape(1, MARGIN + col_w + Inches(0.2),
                                  CONTENT_TOP, Inches(0.02), CONTENT_H)
    div.fill.solid()
    div.fill.fore_color.rgb = ST_LIGHT_GREY
    div.line.fill.background()
    # Right
    right_left = MARGIN + col_w + Inches(0.5)
    _add_box(slide, right_left, CONTENT_TOP, col_w, Inches(0.4),
             cf.get("right_title", ""), font_size=16, bold=True, color=ST_RED)
    _add_box(slide, right_left, CONTENT_TOP + Inches(0.5), col_w, CONTENT_H - Inches(0.5),
             cf.get("right_content", ""), font_size=15, color=ST_DARK)


def _render_three_column(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    col_w = (W - 4 * MARGIN) / 3
    for i, (title_key, content_key) in enumerate([
        ("col1_title", "col1_content"),
        ("col2_title", "col2_content"),
        ("col3_title", "col3_content"),
    ]):
        left = MARGIN + i * (col_w + MARGIN)
        _add_box(slide, left, CONTENT_TOP, col_w, Inches(0.4),
                 cf.get(title_key, ""), font_size=15, bold=True, color=ST_RED)
        _add_box(slide, left, CONTENT_TOP + Inches(0.5), col_w, CONTENT_H - Inches(0.5),
                 cf.get(content_key, ""), font_size=13, color=ST_DARK)


def _render_comparison_2up(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    col_w = (W - 3 * MARGIN) / 2
    for i, (t_key, b_key) in enumerate([("left_title", "left_bullets"), ("right_title", "right_bullets")]):
        left = MARGIN + i * (col_w + MARGIN)
        # Column header with background
        hdr = slide.shapes.add_shape(1, left, CONTENT_TOP, col_w, Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = ST_RED
        hdr.line.fill.background()
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cf.get(t_key, "")
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = ST_WHITE
        # Bullets
        bullets = cf.get(b_key, [])
        content = "\n".join(f"• {b}" for b in bullets)
        _add_box(slide, left, CONTENT_TOP + Inches(0.55), col_w, CONTENT_H - Inches(0.55),
                 content, font_size=14, color=ST_DARK)


def _render_process_flow(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    steps = cf.get("steps", [])
    if not steps:
        return
    step_w = (W - 2 * MARGIN) / max(len(steps), 1)
    arrow_color = ST_RED
    for idx, step in enumerate(steps):
        left = MARGIN + idx * step_w
        box_w = step_w - Inches(0.15)
        # Step number circle (approximated as square)
        num_box = slide.shapes.add_shape(1, left, CONTENT_TOP, Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = ST_RED
        num_box.line.fill.background()
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(idx + 1)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = ST_WHITE
        # Step title
        _add_box(slide, left, CONTENT_TOP + Inches(0.55), box_w, Inches(0.4),
                 step.get("title", ""), font_size=14, bold=True, color=ST_DARK)
        # Step description
        _add_box(slide, left, CONTENT_TOP + Inches(1.05), box_w, CONTENT_H - Inches(1.05),
                 step.get("description", ""), font_size=12, color=ST_GREY)
        # Arrow between steps
        if idx < len(steps) - 1:
            arr_left = left + box_w + Inches(0.02)
            arr = slide.shapes.add_shape(1, arr_left, CONTENT_TOP + Inches(0.1),
                                          Inches(0.12), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = arrow_color
            arr.line.fill.background()


def _render_cards_row(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    cards = cf.get("cards", [])
    if not cards:
        return
    card_w = (W - 2 * MARGIN - Inches(0.2) * (len(cards) - 1)) / max(len(cards), 1)
    for idx, card in enumerate(cards):
        left = MARGIN + idx * (card_w + Inches(0.2))
        # Card box with light grey background
        box = slide.shapes.add_shape(1, left, CONTENT_TOP, card_w, CONTENT_H)
        box.fill.solid()
        box.fill.fore_color.rgb = ST_LIGHT_GREY
        box.line.color.rgb = ST_GREY
        # Card title (red bar at top of card)
        title_box = slide.shapes.add_shape(1, left, CONTENT_TOP, card_w, Inches(0.45))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = ST_RED
        title_box.line.fill.background()
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = card.get("title", "")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = ST_WHITE
        # Card content
        _add_box(slide, left + Inches(0.1), CONTENT_TOP + Inches(0.55),
                 card_w - Inches(0.2), CONTENT_H - Inches(0.65),
                 card.get("content", ""), font_size=12, color=ST_DARK)


def _render_quote_highlight(slide, cf: dict) -> None:
    _message_bar(slide, "")
    # Large quote mark
    _add_box(slide, MARGIN, Inches(1.0), Inches(1.0), Inches(1.2),
             "“", font_size=72, bold=True, color=ST_RED)
    _add_box(slide, MARGIN + Inches(0.8), Inches(2.0), W - 2 * MARGIN - Inches(0.8), Inches(2.5),
             cf.get("quote", ""), font_size=24, bold=True, color=ST_DARK,
             align=PP_ALIGN.CENTER)
    if cf.get("attribution"):
        _add_box(slide, MARGIN, Inches(4.8), W - 2 * MARGIN, Inches(0.5),
                 f"— {cf['attribution']}", font_size=14, color=ST_GREY, align=PP_ALIGN.RIGHT)


def _render_content_placeholder(slide, cf: dict) -> None:
    _message_bar(slide, cf.get("title", ""))
    hint = cf.get("placeholder_hint", "此处放置图表")
    ph = slide.shapes.add_shape(1, MARGIN, CONTENT_TOP, W - 2 * MARGIN, CONTENT_H)
    ph.fill.solid()
    ph.fill.fore_color.rgb = ST_LIGHT_GREY
    ph.line.color.rgb = ST_GREY
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {hint} ]"
    run.font.size = Pt(18)
    run.font.color.rgb = ST_GREY
    run.font.italic = True


def _render_section_divider(slide, cf: dict) -> None:
    # Full red background
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ST_RED
    bg.line.fill.background()
    _add_box(slide, MARGIN, Inches(2.5), W - 2 * MARGIN, Inches(1.8),
             cf.get("title", ""), font_size=36, bold=True,
             color=ST_WHITE, align=PP_ALIGN.CENTER)
    if cf.get("subtitle"):
        _add_box(slide, MARGIN, Inches(4.5), W - 2 * MARGIN, Inches(0.8),
                 cf["subtitle"], font_size=20, color=ST_WHITE, align=PP_ALIGN.CENTER)


_RENDERERS = {
    "title-slide": _render_title_slide,
    "title-bullets": _render_title_bullets,
    "title-image": _render_title_image,
    "two-column": _render_two_column,
    "three-column": _render_three_column,
    "product-comparison-2up": _render_comparison_2up,
    "process-flow": _render_process_flow,
    "cards-row": _render_cards_row,
    "quote-highlight": _render_quote_highlight,
    "content-placeholder": _render_content_placeholder,
    "section-divider": _render_section_divider,
}


def build_pptx(items: list[dict[str, Any]]) -> bytes:
    """Render a list of SlidePlanItem dicts to .pptx bytes."""
    prs = _prs()
    for item in sorted(items, key=lambda x: x.get("order", 0)):
        archetype = item.get("archetype", "title-bullets")
        cf = item.get("content_fields", {})
        slide = _blank_slide(prs)
        renderer = _RENDERERS.get(archetype, _render_title_bullets)
        renderer(slide, cf)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
